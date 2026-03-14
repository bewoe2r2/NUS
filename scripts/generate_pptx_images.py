"""
Bewo Pitch Deck — Image-based PPTX
Embeds high-res React slide screenshots as full-bleed backgrounds.
"""

from pptx import Presentation
from pptx.util import Inches, Emu
import os

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank

captures_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slide_captures")

for i in range(1, 7):
    img_path = os.path.join(captures_dir, f"slide_{i}.png")
    if not os.path.exists(img_path):
        print(f"Missing: {img_path}")
        continue

    slide = prs.slides.add_slide(blank_layout)
    # Full-bleed image covering entire slide
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    print(f"Added slide {i}")

output_path = os.path.join(os.path.dirname(captures_dir), "Bewo_Pitch_Deck_v2.pptx")
prs.save(output_path)
print(f"\nSaved: {output_path}")
