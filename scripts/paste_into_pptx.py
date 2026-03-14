"""
Paste captured slide PNGs into a clean 16:9 PPTX.
One full-bleed image per slide — pixel-perfect from React rendering.
"""
import os
from pptx import Presentation
from pptx.util import Inches

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NUM_SLIDES = 10

captures_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slide_captures")
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "Bewo_Pitch_v7.pptx")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

for i in range(1, NUM_SLIDES + 1):
    img_path = os.path.join(captures_dir, f"slide_{i}.png")
    if not os.path.exists(img_path):
        print(f"WARNING: {img_path} not found, skipping")
        continue
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    print(f"  Added slide {i}")

prs.save(output_path)
print(f"\nSaved: {output_path}")
