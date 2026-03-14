"""
Bewo Pitch Deck — NUS Synapxe IMDA AI Innovation Challenge 2026
Editorial white-paper design inspired by analyst-portfolio aesthetic.
10-slide narrative arc: Hook → Problem → Solution → Architecture → Patient → Nurse → AI → Safety → Impact → Close
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math
import os

# ═══════════════════════════════════════════════════════════
# DESIGN SYSTEM — Wall Street Paper (analyst-portfolio)
# ═══════════════════════════════════════════════════════════

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Backgrounds
PAPER = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE_BG = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_BLUE_BG = RGBColor(0xEE, 0xF2, 0xFF)

# Text
INK = RGBColor(0x1A, 0x1A, 0x2E)
SECONDARY = RGBColor(0x5A, 0x5A, 0x7A)
MUTED = RGBColor(0x8A, 0x8A, 0xA0)

# Accents
NAVY = RGBColor(0x2B, 0x4C, 0x8C)
EMERALD = RGBColor(0x15, 0x80, 0x3D)
CRIMSON = RGBColor(0xB9, 0x1C, 0x1C)
AMBER = RGBColor(0xD9, 0x77, 0x06)

# Borders
HAIRLINE = RGBColor(0xE5, 0xE5, 0xE5)
STRONG = RGBColor(0x1A, 0x1A, 0x2E)

# Fonts
SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

# Spacing
MARGIN_L = Inches(1.0)
MARGIN_R = Inches(1.0)
CONTENT_W = Inches(11.333)

# ═══════════════════════════════════════════════════════════
# PRESENTATION SETUP
# ═══════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════

def add_bg(slide, color=PAPER):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0), corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=12, color=INK, bold=False, font_name=SANS, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_para(tf, text, font_size=12, color=INK, bold=False, font_name=SANS, alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def mixed_text(slide, left, top, width, height, runs_data):
    """Add text box with multiple styled runs in a single paragraph.
    runs_data: list of (text, font_size, color, bold, font_name) tuples"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    for text, font_size, color, bold, font_name in runs_data:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txbox


def add_divider(slide, left, top, width, color=HAIRLINE, thickness=Pt(1)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_accent_card(slide, left, top, width, height, accent_color=NAVY, fill=WHITE):
    """Card with left accent border (MetricCard style)."""
    # Main card bg
    card = add_shape(slide, left, top, width, height, fill_color=fill, border_color=HAIRLINE, border_width=Pt(1))
    # Left accent bar
    add_shape(slide, left, top, Pt(3), height, fill_color=accent_color)
    return card


def add_metric_card(slide, x, y, label, value, sub, accent_color=NAVY):
    """MetricCard: left-border accent, label/value/sub stack."""
    card_w = Inches(2.4)
    card_h = Inches(1.3)
    add_accent_card(slide, x, y, card_w, card_h, accent_color=accent_color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), card_w - Inches(0.3), Inches(0.2),
             label.upper(), font_size=9, color=MUTED, font_name=MONO)
    add_text(slide, x + Inches(0.2), y + Inches(0.4), card_w - Inches(0.3), Inches(0.45),
             value, font_size=28, color=INK, bold=True, font_name=MONO)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), card_w - Inches(0.3), Inches(0.25),
             sub, font_size=10, color=SECONDARY)


def add_bar(slide, left, top, width, height, fill_pct, bg_color, bar_color):
    add_shape(slide, left, top, width, height, fill_color=bg_color, corner_radius=True)
    bar_w = int(width * fill_pct / 100)
    if bar_w > 0:
        add_shape(slide, left, top, bar_w, height, fill_color=bar_color, corner_radius=True)


def add_dot(slide, x, y, color, size=Inches(0.1)):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot


def add_label_tag(slide, x, y, text, color=NAVY, bg=LIGHT_BLUE_BG):
    """Small pill tag."""
    tag_w = Inches(0.15 + len(text) * 0.065)
    tag_h = Inches(0.25)
    add_shape(slide, x, y, tag_w, tag_h, fill_color=bg, corner_radius=True)
    add_text(slide, x + Inches(0.08), y + Inches(0.02), tag_w - Inches(0.1), tag_h,
             text, font_size=8, color=color, bold=True, font_name=MONO)
    return tag_w


def slide_footer(slide, text="NUS x Synapxe x IMDA  |  AI Innovation Challenge 2026"):
    add_text(slide, MARGIN_L, Inches(7.05), Inches(8), Inches(0.3),
             text, font_size=8, color=MUTED, font_name=MONO)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: HOOK — "440,000 diabetics. Zero early warning."
# ═══════════════════════════════════════════════════════════

def build_slide_1():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    # Top-left accent label
    add_text(slide, MARGIN_L, Inches(0.8), Inches(4), Inches(0.3),
             "SINGAPORE'S SILENT CRISIS", font_size=9, color=CRIMSON, font_name=MONO)

    # Giant headline with mixed colors
    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(1.3), Inches(7), Inches(3.5))
    tf = txbox.text_frame
    tf.word_wrap = True

    # "440,000" in crimson
    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = "440,000"
    r.font.size = Pt(72)
    r.font.color.rgb = CRIMSON
    r.font.bold = True
    r.font.name = SERIF

    # "diabetics" in muted
    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = "diabetics"
    r.font.size = Pt(72)
    r.font.color.rgb = MUTED
    r.font.bold = True
    r.font.name = SERIF

    # "zero early warning." — "zero" in crimson, rest in ink
    p3 = tf.add_paragraph()
    r1 = p3.add_run()
    r1.text = "zero "
    r1.font.size = Pt(72)
    r1.font.color.rgb = CRIMSON
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p3.add_run()
    r2.text = "early warning."
    r2.font.size = Pt(72)
    r2.font.color.rgb = INK
    r2.font.bold = True
    r2.font.name = SERIF

    # Subtitle
    add_text(slide, MARGIN_L, Inches(5.0), Inches(6), Inches(0.6),
             "61% of diabetes spending goes to preventable hospitalizations.",
             font_size=15, color=SECONDARY)

    # Divider
    add_divider(slide, MARGIN_L, Inches(5.7), Inches(5.5))

    # Right side: 4 stat cards (MetricCard style)
    rx = Inches(8.5)
    stats = [
        ("ER Visit Cost", "$8,800", "Per admission avg", CRIMSON),
        ("Annual Spend", "$4.0B", "Singapore diabetes", CRIMSON),
        ("Preventable", "40%", "ER admissions avoidable", AMBER),
        ("Early Detection", "0 hrs", "Between appointments", RGBColor(0x64, 0x64, 0x80)),
    ]
    for i, (label, value, sub, color) in enumerate(stats):
        y = Inches(1.2) + Inches(i * 1.45)
        add_metric_card(slide, rx, y, label, value, sub, accent_color=color)

    # Footer
    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════

def build_slide_2():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.8), Inches(5), Inches(0.3),
             "THE PROBLEM", font_size=9, color=CRIMSON, font_name=MONO)

    # Headline
    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(1.2), Inches(10), Inches(1.0))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "The system demands patients "
    r1.font.size = Pt(36)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "manage their disease."
    r2.font.size = Pt(36)
    r2.font.color.rgb = CRIMSON
    r2.font.bold = True
    r2.font.name = SERIF

    # 3 problem cards
    problems = [
        ("Manual Logging", "Patients expected to track glucose, meals, medication daily. Most don't.", "The average diabetic is asked to log 4-6 data points per day — indefinitely."),
        ("Reactive Care", "Clinicians only see patients every 3-6 months. Crises happen in hours.", "Between appointments, patients are on their own. No monitoring. No safety net."),
        ("Non-Compliance Blame", "When patients fail, the system calls it 'non-adherence.'", "The problem isn't non-compliance — it's a system designed around clinic visits, not daily life."),
    ]

    card_w = Inches(3.4)
    card_h = Inches(2.8)
    card_y = Inches(2.6)

    for i, (title, desc, detail) in enumerate(problems):
        cx = MARGIN_L + Inches(i * 3.7)

        # Card with left accent
        add_accent_card(slide, cx, card_y, card_w, card_h, accent_color=CRIMSON, fill=WHITE)

        # Number
        add_text(slide, cx + Inches(0.25), card_y + Inches(0.2), Inches(0.5), Inches(0.4),
                 f"0{i+1}", font_size=24, color=RGBColor(0xE0, 0xE0, 0xE8), bold=True, font_name=MONO)

        # Title
        add_text(slide, cx + Inches(0.25), card_y + Inches(0.65), card_w - Inches(0.5), Inches(0.35),
                 title, font_size=18, color=INK, bold=True, font_name=SERIF)

        # Description
        add_text(slide, cx + Inches(0.25), card_y + Inches(1.05), card_w - Inches(0.5), Inches(0.7),
                 desc, font_size=12, color=SECONDARY)

        # Detail
        add_divider(slide, cx + Inches(0.25), card_y + Inches(1.85), card_w - Inches(0.5))
        add_text(slide, cx + Inches(0.25), card_y + Inches(2.0), card_w - Inches(0.5), Inches(0.6),
                 detail, font_size=10, color=MUTED, font_name=SANS)

    # Bottom insight
    add_shape(slide, MARGIN_L, Inches(5.8), CONTENT_W, Inches(0.6), fill_color=LIGHT_BLUE_BG)
    add_text(slide, MARGIN_L + Inches(0.3), Inches(5.9), Inches(10), Inches(0.4),
             "What if the system watched — so the patient didn't have to?",
             font_size=14, color=NAVY, bold=True, font_name=SERIF)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: BEWO — THE SOLUTION
# ═══════════════════════════════════════════════════════════

def build_slide_3():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "OUR SOLUTION", font_size=9, color=NAVY, font_name=MONO)

    # Headline
    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(10), Inches(0.8))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Bewo predicts crises "
    r1.font.size = Pt(36)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "48 hours early."
    r2.font.size = Pt(36)
    r2.font.color.rgb = NAVY
    r2.font.bold = True
    r2.font.name = SERIF

    # Tagline
    add_text(slide, MARGIN_L, Inches(1.75), Inches(8), Inches(0.4),
             "The patient's job is to live their life, not to manage their disease.",
             font_size=14, color=SECONDARY, font_name=SERIF)

    add_divider(slide, MARGIN_L, Inches(2.3), CONTENT_W)

    # 3 pillar cards
    pillars = [
        {
            "title": "HMM Engine",
            "subtitle": "PREDICTIVE INTELLIGENCE",
            "color": AMBER,
            "bullets": [
                "3-state Viterbi: STABLE / WARNING / CRISIS",
                "Baum-Welch EM learns per-patient transitions",
                "9 biomarkers tracked passively",
            ],
        },
        {
            "title": "Agentic AI",
            "subtitle": "AUTONOMOUS ACTIONS",
            "color": NAVY,
            "bullets": [
                "18 clinical tools (book, alert, adjust, escalate)",
                "5-turn ReAct loop: Observe / Think / Act",
                "Cross-session memory with Gemini extraction",
            ],
        },
        {
            "title": "Behavioral Science",
            "subtitle": "SUSTAINED ENGAGEMENT",
            "color": EMERALD,
            "bullets": [
                "$5/week loss-aversion voucher (Prospect Theory)",
                "7-day streaks with CDC-aligned gamification",
                "Mood-aware adaptive Singlish tone",
            ],
        },
    ]

    card_w = Inches(3.5)
    card_h = Inches(3.8)
    card_y = Inches(2.7)

    for i, pillar in enumerate(pillars):
        cx = MARGIN_L + Inches(i * 3.8)
        color = pillar["color"]

        # Card
        card = add_shape(slide, cx, card_y, card_w, card_h, fill_color=WHITE,
                         border_color=HAIRLINE, border_width=Pt(1))

        # Top accent bar
        add_shape(slide, cx, card_y, card_w, Pt(4), fill_color=color)

        # Subtitle
        add_text(slide, cx + Inches(0.25), card_y + Inches(0.25), card_w - Inches(0.5), Inches(0.25),
                 pillar["subtitle"], font_size=9, color=color, bold=True, font_name=MONO)

        # Title
        add_text(slide, cx + Inches(0.25), card_y + Inches(0.55), card_w - Inches(0.5), Inches(0.4),
                 pillar["title"], font_size=22, color=INK, bold=True, font_name=SERIF)

        # Divider
        add_divider(slide, cx + Inches(0.25), card_y + Inches(1.05), card_w - Inches(0.5))

        # Bullets
        for bi, bullet in enumerate(pillar["bullets"]):
            by = card_y + Inches(1.25 + bi * 0.7)
            add_dot(slide, cx + Inches(0.3), by + Inches(0.05), color, size=Inches(0.08))
            add_text(slide, cx + Inches(0.5), by, card_w - Inches(0.7), Inches(0.55),
                     bullet, font_size=12, color=SECONDARY)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: HOW IT WORKS (Architecture)
# ═══════════════════════════════════════════════════════════

def build_slide_4():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "ARCHITECTURE", font_size=9, color=NAVY, font_name=MONO)

    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(6), Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Real AI. "
    r1.font.size = Pt(32)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "Not a wrapper."
    r2.font.size = Pt(32)
    r2.font.color.rgb = NAVY
    r2.font.bold = True
    r2.font.name = SERIF

    add_divider(slide, MARGIN_L, Inches(1.65), Inches(7.5))

    # Pipeline steps (left column)
    steps = [
        ("1", "Patient Data", "Passive sensors — glucose, BP, heart rate, SpO2, weight, steps, sleep, mood, HbA1c", MUTED),
        ("2", "HMM Engine", "Viterbi + Baum-Welch (9 biomarkers -> STABLE / WARNING / CRISIS)", AMBER),
        ("3", "Merlion Risk", "A*STAR ARIMA: 48h glucose forecast + Monte Carlo confidence intervals", RGBColor(0x1E, 0x60, 0xA0)),
        ("4", "Gemini Agent", "5-turn ReAct with 18 tools + cross-session episodic memory", NAVY),
        ("5", "Safety Filter", "6-dimension check: medical claims, hallucination, cultural, scope, emotional, dangerous", CRIMSON),
        ("6", "SEA-LION", "AI Singapore: Singlish translation + cultural nuance for elderly patients", EMERALD),
    ]

    step_x = MARGIN_L
    step_w = Inches(7.0)
    step_h = Inches(0.7)

    for i, (num, title, desc, color) in enumerate(steps):
        sy = Inches(1.95) + Inches(i * 0.82)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, step_x, sy + Inches(0.05), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(slide, step_x + Inches(0.02), sy + Inches(0.07), Inches(0.35), Inches(0.35),
                 num, font_size=14, color=WHITE, bold=True, font_name=MONO, alignment=PP_ALIGN.CENTER)

        # Title
        add_text(slide, step_x + Inches(0.5), sy, Inches(2), Inches(0.3),
                 title, font_size=14, color=INK, bold=True, font_name=SERIF)

        # Description
        add_text(slide, step_x + Inches(0.5), sy + Inches(0.3), step_w - Inches(0.8), Inches(0.35),
                 desc, font_size=10, color=SECONDARY)

        # Connector line (except last)
        if i < len(steps) - 1:
            add_shape(slide, step_x + Inches(0.16), sy + Inches(0.45), Pt(2), Inches(0.38), fill_color=HAIRLINE)

    # Right side: Key numbers
    rx = Inches(8.8)
    add_text(slide, rx, Inches(1.95), Inches(3.5), Inches(0.3),
             "KEY NUMBERS", font_size=9, color=NAVY, font_name=MONO)

    add_divider(slide, rx, Inches(2.3), Inches(3.5))

    key_nums = [
        ("9", "biomarkers tracked"),
        ("18", "autonomous tools"),
        ("6", "safety dimensions"),
        ("48h", "forecast horizon"),
        ("5", "ReAct turns max"),
        ("16+", "drug interaction pairs"),
        ("39", "drug-class mappings"),
        ("3", "HMM states"),
    ]

    for i, (num, label) in enumerate(key_nums):
        ny = Inches(2.5) + Inches(i * 0.55)
        add_text(slide, rx, ny, Inches(1.0), Inches(0.35),
                 num, font_size=20, color=NAVY, bold=True, font_name=MONO)
        add_text(slide, rx + Inches(1.0), ny + Inches(0.05), Inches(2.5), Inches(0.3),
                 label, font_size=12, color=SECONDARY)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: PATIENT EXPERIENCE
# ═══════════════════════════════════════════════════════════

def build_slide_5():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "PATIENT EXPERIENCE", font_size=9, color=EMERALD, font_name=MONO)

    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(8), Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Uncle Tan "
    r1.font.size = Pt(32)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "never opens an app."
    r2.font.size = Pt(32)
    r2.font.color.rgb = EMERALD
    r2.font.bold = True
    r2.font.name = SERIF

    add_divider(slide, MARGIN_L, Inches(1.6), CONTENT_W)

    # Left: Phone mockup
    px = MARGIN_L + Inches(0.5)
    py = Inches(1.9)
    pw = Inches(3.5)
    ph = Inches(5.0)

    # Phone outer frame
    phone = add_shape(slide, px, py, pw, ph, fill_color=WHITE,
                      border_color=STRONG, border_width=Pt(2), corner_radius=True)

    # Phone header bar
    add_shape(slide, px, py, pw, Inches(0.5), fill_color=NAVY, corner_radius=True)
    add_text(slide, px + Inches(0.2), py + Inches(0.1), pw - Inches(0.4), Inches(0.3),
             "Bewo", font_size=14, color=WHITE, bold=True, font_name=SERIF)

    # Chat messages
    messages = [
        (False, "Uncle ah, your sugar going up leh. Can drink water first?"),
        (True, "Ok lah"),
        (False, "Good! One more day steady and you hit 7-day streak \u2014 $5 voucher!"),
        (True, "Wah shiok, thanks ah"),
        (False, "No problem uncle! Your trend looking steady this week. Keep it up!"),
    ]

    my = py + Inches(0.65)
    for is_user, msg in messages:
        # Sender label
        sender = "Uncle Tan" if is_user else "Bewo"
        sender_color = MUTED if is_user else NAVY
        add_text(slide, px + Inches(0.2), my, Inches(1.5), Inches(0.2),
                 sender, font_size=8, color=sender_color, bold=True, font_name=MONO)
        my += Inches(0.2)

        # Bubble
        lines = math.ceil(len(msg) / 40)
        bh = Inches(0.18 + lines * 0.2)
        bw = pw - Inches(0.4)

        if is_user:
            bubble_bg = SUBTLE_BG
            text_color = INK
        else:
            bubble_bg = LIGHT_BLUE_BG
            text_color = NAVY

        add_shape(slide, px + Inches(0.2), my, bw, bh, fill_color=bubble_bg, corner_radius=True)
        add_text(slide, px + Inches(0.3), my + Inches(0.04), bw - Inches(0.2), bh,
                 msg, font_size=10, color=text_color)
        my += bh + Inches(0.1)

    # Right: Feature bullets
    fx = Inches(5.8)
    fy = Inches(2.0)

    features = [
        ("Zero-effort monitoring", "Passive sensors collect data. Patient does nothing.", EMERALD),
        ("Proactive check-ins", "AI initiates contact — patient doesn't need to remember.", NAVY),
        ("Singlish-native", "SEA-LION translates to culturally appropriate language.", AMBER),
        ("Loss-aversion voucher", "$5/week gamification based on Prospect Theory.", EMERALD),
        ("Mood-aware tone", "Detects frustration, adjusts language. Never patronizing.", NAVY),
    ]

    for i, (title, desc, color) in enumerate(features):
        fy_i = fy + Inches(i * 1.0)

        # Accent card
        add_accent_card(slide, fx, fy_i, Inches(6.2), Inches(0.8), accent_color=color, fill=WHITE)

        add_text(slide, fx + Inches(0.25), fy_i + Inches(0.1), Inches(5.5), Inches(0.3),
                 title, font_size=14, color=INK, bold=True, font_name=SERIF)
        add_text(slide, fx + Inches(0.25), fy_i + Inches(0.42), Inches(5.5), Inches(0.3),
                 desc, font_size=11, color=SECONDARY)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: NURSE DASHBOARD
# ═══════════════════════════════════════════════════════════

def build_slide_6():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "CLINICAL INTERFACE", font_size=9, color=NAVY, font_name=MONO)

    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(8), Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "1 nurse. 100+ patients. "
    r1.font.size = Pt(32)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "Instant triage."
    r2.font.size = Pt(32)
    r2.font.color.rgb = NAVY
    r2.font.bold = True
    r2.font.name = SERIF

    add_divider(slide, MARGIN_L, Inches(1.6), CONTENT_W)

    # Left: feature list
    features = [
        "Auto-SBAR reports (Situation / Background / Assessment / Recommendation)",
        "Multi-patient triage scoring with urgency classification",
        "Drug interaction checking (16+ pairs, severity-ranked)",
        "Baum-Welch per-patient HMM parameter learning",
        "Caregiver burden tracking (4-factor scoring, 0-100)",
    ]

    for i, feat in enumerate(features):
        fy = Inches(1.9) + Inches(i * 0.42)
        add_dot(slide, MARGIN_L, fy + Inches(0.06), NAVY, size=Inches(0.08))
        add_text(slide, MARGIN_L + Inches(0.2), fy, Inches(4.2), Inches(0.35),
                 feat, font_size=11, color=SECONDARY)

    # Right: Dashboard mock (triage table)
    dx = Inches(5.5)
    dy = Inches(1.9)
    dw = Inches(7.0)
    dh = Inches(4.8)

    # Dashboard frame
    add_shape(slide, dx, dy, dw, dh, fill_color=WHITE, border_color=STRONG, border_width=Pt(1.5))

    # Dashboard header
    add_shape(slide, dx, dy, dw, Inches(0.45), fill_color=NAVY)
    add_text(slide, dx + Inches(0.2), dy + Inches(0.08), Inches(2), Inches(0.3),
             "Nurse Triage Dashboard", font_size=11, color=WHITE, bold=True, font_name=SANS)
    add_text(slide, dx + dw - Inches(2.5), dy + Inches(0.1), Inches(2.2), Inches(0.25),
             "3 patients need attention", font_size=9, color=RGBColor(0xBB, 0xCC, 0xFF), alignment=PP_ALIGN.RIGHT)

    # Column headers
    hy = dy + Inches(0.55)
    add_text(slide, dx + Inches(0.2), hy, Inches(2), Inches(0.25),
             "PATIENT", font_size=8, color=MUTED, bold=True, font_name=MONO)
    add_text(slide, dx + Inches(2.4), hy, Inches(0.8), Inches(0.25),
             "STATE", font_size=8, color=MUTED, bold=True, font_name=MONO)
    add_text(slide, dx + Inches(3.3), hy, Inches(2.5), Inches(0.25),
             "CLINICAL NOTE", font_size=8, color=MUTED, bold=True, font_name=MONO)
    add_text(slide, dx + Inches(6.0), hy, Inches(0.8), Inches(0.25),
             "URGENCY", font_size=8, color=MUTED, bold=True, font_name=MONO)

    add_divider(slide, dx + Inches(0.1), hy + Inches(0.25), dw - Inches(0.2), color=STRONG)

    # Patient rows
    patients = [
        ("Tan Ah Kow, 72", "CRISIS", CRIMSON, "Glucose 15.2 — rising 6h straight", "IMMEDIATE", RGBColor(0xFF, 0xF0, 0xF0)),
        ("Lim Mei Hua, 68", "WARNING", AMBER, "Missed 2 doses, HbA1c trending up", "SOON", RGBColor(0xFF, 0xFA, 0xF0)),
        ("Wong Siew Lan, 75", "WARNING", AMBER, "Logging gap 8h, BP elevated", "SOON", RGBColor(0xFF, 0xFA, 0xF0)),
        ("Chen Da Wei, 61", "STABLE", EMERALD, "All vitals normal, 14-day streak", "MONITOR", WHITE),
        ("Lee Bee Hoon, 79", "STABLE", EMERALD, "Glucose stable 5.8-6.4 range", "STABLE", WHITE),
    ]

    for i, (name, state, color, note, urgency, row_bg) in enumerate(patients):
        ry = hy + Inches(0.35 + i * 0.7)

        # Row background
        add_shape(slide, dx + Inches(0.1), ry, dw - Inches(0.2), Inches(0.6), fill_color=row_bg)

        # Name
        add_text(slide, dx + Inches(0.2), ry + Inches(0.1), Inches(2), Inches(0.2),
                 name, font_size=11, color=INK, bold=True)

        # State badge
        badge_w = Inches(0.75)
        add_shape(slide, dx + Inches(2.4), ry + Inches(0.12), badge_w, Inches(0.22),
                  fill_color=color, corner_radius=True)
        add_text(slide, dx + Inches(2.4), ry + Inches(0.12), badge_w, Inches(0.22),
                 state, font_size=7, color=WHITE, bold=True, font_name=MONO, alignment=PP_ALIGN.CENTER)

        # Note
        add_text(slide, dx + Inches(3.3), ry + Inches(0.1), Inches(2.5), Inches(0.4),
                 note, font_size=9, color=SECONDARY)

        # Urgency
        add_text(slide, dx + Inches(6.0), ry + Inches(0.15), Inches(0.9), Inches(0.2),
                 urgency, font_size=8, color=color, bold=True, font_name=MONO)

        # Row divider
        if i < len(patients) - 1:
            add_divider(slide, dx + Inches(0.1), ry + Inches(0.62), dw - Inches(0.2))

    # SBAR sample below features
    sbar_y = Inches(4.2)
    add_text(slide, MARGIN_L, sbar_y, Inches(4.2), Inches(0.25),
             "SAMPLE SBAR REPORT", font_size=9, color=NAVY, font_name=MONO)

    sbar_items = [
        ("S:", "Glucose 15.2 mmol/L, rising trend 6 hours"),
        ("B:", "Type 2 DM, HbA1c 8.1%, metformin + glipizide"),
        ("A:", "HMM state: CRISIS. Merlion 48h forecast: continued rise"),
        ("R:", "Urgent review. Consider insulin adjustment. Alert caregiver."),
    ]

    for i, (letter, text) in enumerate(sbar_items):
        sy = sbar_y + Inches(0.3 + i * 0.35)
        add_text(slide, MARGIN_L, sy, Inches(0.3), Inches(0.25),
                 letter, font_size=10, color=NAVY, bold=True, font_name=MONO)
        add_text(slide, MARGIN_L + Inches(0.3), sy, Inches(3.8), Inches(0.3),
                 text, font_size=10, color=SECONDARY)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 7: AI INTELLIGENCE
# ═══════════════════════════════════════════════════════════

def build_slide_7():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "AI INTELLIGENCE", font_size=9, color=NAVY, font_name=MONO)

    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(10), Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "The agent that "
    r1.font.size = Pt(32)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "learns and remembers."
    r2.font.size = Pt(32)
    r2.font.color.rgb = NAVY
    r2.font.bold = True
    r2.font.name = SERIF

    add_divider(slide, MARGIN_L, Inches(1.6), CONTENT_W)

    # 2x2 capability grid
    capabilities = [
        {
            "title": "Cross-Session Memory",
            "desc": "Remembers 'metformin makes you nauseous.' Adjusts all future advice. Episodic + semantic + preference memory types.",
            "color": NAVY,
            "icon": "MEM",
        },
        {
            "title": "Outcome-Based Tool Selection",
            "desc": "Learns which tools work for EACH patient. 14-day half-life decay. Tool effectiveness scored per HMM state.",
            "color": AMBER,
            "icon": "OPT",
        },
        {
            "title": "Proactive Scheduler",
            "desc": "6 triggers: glucose rising, sustained risk, logging gap, medication nudge, streak save, mood follow-up.",
            "color": EMERALD,
            "icon": "PRO",
        },
        {
            "title": "Drug Interaction Checker",
            "desc": "16 interaction pairs, 39 drug-class mappings. Severity-ranked. Auto-blocks CONTRAINDICATED combinations.",
            "color": CRIMSON,
            "icon": "RX",
        },
    ]

    card_w = Inches(5.3)
    card_h = Inches(2.0)

    for i, cap in enumerate(capabilities):
        col = i % 2
        row = i // 2
        cx = MARGIN_L + Inches(col * 5.7)
        cy = Inches(1.9) + Inches(row * 2.3)
        color = cap["color"]

        # Card
        add_accent_card(slide, cx, cy, card_w, card_h, accent_color=color, fill=WHITE)

        # Icon badge
        add_shape(slide, cx + Inches(0.2), cy + Inches(0.2), Inches(0.5), Inches(0.35),
                  fill_color=color, corner_radius=True)
        add_text(slide, cx + Inches(0.2), cy + Inches(0.23), Inches(0.5), Inches(0.3),
                 cap["icon"], font_size=9, color=WHITE, bold=True, font_name=MONO, alignment=PP_ALIGN.CENTER)

        # Title
        add_text(slide, cx + Inches(0.85), cy + Inches(0.2), card_w - Inches(1.1), Inches(0.3),
                 cap["title"], font_size=16, color=INK, bold=True, font_name=SERIF)

        # Description
        add_text(slide, cx + Inches(0.25), cy + Inches(0.7), card_w - Inches(0.5), Inches(1.0),
                 cap["desc"], font_size=12, color=SECONDARY)

    # Center stat
    add_shape(slide, Inches(3.5), Inches(6.3), Inches(6.5), Inches(0.5), fill_color=LIGHT_BLUE_BG)
    add_text(slide, Inches(3.5), Inches(6.35), Inches(6.5), Inches(0.4),
             "18 tools  x  3 states  x  per-patient learning  =  personalized care at scale",
             font_size=12, color=NAVY, bold=True, font_name=MONO, alignment=PP_ALIGN.CENTER)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: SAFETY & TRUST
# ═══════════════════════════════════════════════════════════

def build_slide_8():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "SAFETY & TRUST", font_size=9, color=CRIMSON, font_name=MONO)

    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(10), Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Healthcare AI demands "
    r1.font.size = Pt(32)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "guardrails."
    r2.font.size = Pt(32)
    r2.font.color.rgb = CRIMSON
    r2.font.bold = True
    r2.font.name = SERIF

    add_divider(slide, MARGIN_L, Inches(1.6), CONTENT_W)

    # 6-dimension safety classifier
    dimensions = [
        ("Medical Claims", "Dosage/diagnosis deferred to doctors. AI suggests, clinicians decide.", EMERALD),
        ("Emotional Mismatch", "Tone matches patient state. Never cheerful during distress.", NAVY),
        ("Hallucination", "Only references real medications and verified clinical data.", CRIMSON),
        ("Cultural Sensitivity", "Elderly Singaporean norms. Respectful address. No assumptions.", AMBER),
        ("Scope Violation", "Stays within chronic disease management. No psychiatric or surgical advice.", NAVY),
        ("Dangerous Advice", "Never suggests stopping medication. Always escalates uncertainty.", CRIMSON),
    ]

    for i, (title, desc, color) in enumerate(dimensions):
        row = i // 2
        col = i % 2
        cx = MARGIN_L + Inches(col * 5.7)
        cy = Inches(1.9) + Inches(row * 1.3)

        # Check mark
        check_bg = add_shape(slide, cx, cy + Inches(0.05), Inches(0.3), Inches(0.3),
                              fill_color=EMERALD, corner_radius=True)
        add_text(slide, cx + Inches(0.02), cy + Inches(0.03), Inches(0.3), Inches(0.3),
                 "\u2713", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Content
        add_text(slide, cx + Inches(0.45), cy, Inches(4.8), Inches(0.3),
                 title, font_size=14, color=INK, bold=True, font_name=SERIF)
        add_text(slide, cx + Inches(0.45), cy + Inches(0.35), Inches(4.8), Inches(0.6),
                 desc, font_size=11, color=SECONDARY)

    # Doctor-gated banner
    banner_y = Inches(5.0)
    add_shape(slide, MARGIN_L, banner_y, CONTENT_W, Inches(0.6), fill_color=RGBColor(0xFF, 0xF0, 0xF0),
              border_color=CRIMSON, border_width=Pt(1))
    add_text(slide, MARGIN_L + Inches(0.3), banner_y + Inches(0.1), Inches(10), Inches(0.4),
             "Doctor-gated: Every medication suggestion requires clinical review before reaching the patient.",
             font_size=13, color=CRIMSON, bold=True)

    # Singapore alignment badges
    badge_y = Inches(5.9)
    add_text(slide, MARGIN_L, badge_y, Inches(1.5), Inches(0.3),
             "ALIGNED WITH:", font_size=9, color=MUTED, font_name=MONO)

    badges = ["Healthier SG", "National AI 2.0", "MOH Standards", "SBAR Protocol"]
    bx = MARGIN_L + Inches(1.6)
    for badge_text in badges:
        tw = add_label_tag(slide, bx, badge_y + Inches(0.02), badge_text, color=NAVY, bg=LIGHT_BLUE_BG)
        bx += tw + Inches(0.15)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: IMPACT & FEASIBILITY
# ═══════════════════════════════════════════════════════════

def build_slide_9():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text(slide, MARGIN_L, Inches(0.6), Inches(4), Inches(0.3),
             "IMPACT & FEASIBILITY", font_size=9, color=EMERALD, font_name=MONO)

    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(8), Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Why it works "
    r1.font.size = Pt(32)
    r1.font.color.rgb = INK
    r1.font.bold = True
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "— and why now."
    r2.font.size = Pt(32)
    r2.font.color.rgb = EMERALD
    r2.font.bold = True
    r2.font.name = SERIF

    add_divider(slide, MARGIN_L, Inches(1.6), CONTENT_W)

    # Left: 4 big metric cards
    metrics = [
        ("48h", "Advance Warning", "Predict crises before they happen", NAVY),
        ("~40%", "ER Reduction", "Preventable admissions avoided", EMERALD),
        ("1:100+", "Nurse Ratio", "vs current 1:20 manual ratio", NAVY),
        ("$0.40", "Server Cost", "Per-patient per-month (Gemini API)", EMERALD),
    ]

    for i, (value, label, sub, color) in enumerate(metrics):
        row = i // 2
        col = i % 2
        mx = MARGIN_L + Inches(col * 2.7)
        my = Inches(1.9) + Inches(row * 1.6)
        add_metric_card(slide, mx, my, label, value, sub, accent_color=color)

    # Right: Business model
    rx = Inches(6.5)
    add_text(slide, rx, Inches(1.9), Inches(5.5), Inches(0.3),
             "BUSINESS MODEL", font_size=9, color=NAVY, font_name=MONO)
    add_divider(slide, rx, Inches(2.25), Inches(5.5))

    # Pricing card
    add_accent_card(slide, rx, Inches(2.4), Inches(5.5), Inches(1.2), accent_color=NAVY, fill=WHITE)
    add_text(slide, rx + Inches(0.25), Inches(2.55), Inches(3), Inches(0.3),
             "B2G SaaS: $3/patient/month", font_size=16, color=INK, bold=True, font_name=SERIF)
    add_text(slide, rx + Inches(0.25), Inches(2.9), Inches(4.5), Inches(0.5),
             "Government pays. Patients use for free. 87% gross margin.",
             font_size=12, color=SECONDARY)

    # Unit economics
    add_text(slide, rx, Inches(3.8), Inches(5.5), Inches(0.3),
             "UNIT ECONOMICS", font_size=9, color=NAVY, font_name=MONO)

    econ = [
        ("1 prevented ER visit", "$8,800 saved", EMERALD),
        ("Cost to prevent (2,900 pt-months)", "$8,700", SECONDARY),
        ("ROI per prevented admission", "+$100 net", EMERALD),
    ]

    for i, (label, value, color) in enumerate(econ):
        ey = Inches(4.15) + Inches(i * 0.35)
        add_text(slide, rx + Inches(0.1), ey, Inches(3.5), Inches(0.25),
                 label, font_size=11, color=SECONDARY)
        add_text(slide, rx + Inches(3.8), ey, Inches(1.6), Inches(0.25),
                 value, font_size=11, color=color, bold=True, font_name=MONO, alignment=PP_ALIGN.RIGHT)

    # Market expansion
    add_text(slide, rx, Inches(5.3), Inches(5.5), Inches(0.3),
             "MARKET", font_size=9, color=NAVY, font_name=MONO)

    markets = [
        ("Singapore", "440K diabetics, 23 polyclinics"),
        ("ASEAN", "56M diabetics, 10 nations"),
        ("Global", "537M diabetics (IDF 2021)"),
    ]

    for i, (region, detail) in enumerate(markets):
        my = Inches(5.65) + Inches(i * 0.32)
        add_text(slide, rx + Inches(0.1), my, Inches(1.2), Inches(0.25),
                 region, font_size=11, color=INK, bold=True, font_name=MONO)
        add_text(slide, rx + Inches(1.5), my, Inches(3.5), Inches(0.25),
                 detail, font_size=11, color=SECONDARY)

    # Bottom: No research breakthroughs needed
    add_shape(slide, MARGIN_L, Inches(6.5), CONTENT_W, Inches(0.5), fill_color=LIGHT_BLUE_BG)
    add_text(slide, MARGIN_L + Inches(0.3), Inches(6.55), Inches(10), Inches(0.4),
             "No research breakthroughs needed. Every component is proven technology.",
             font_size=13, color=NAVY, bold=True, font_name=SERIF)

    slide_footer(slide)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: CLOSE
# ═══════════════════════════════════════════════════════════

def build_slide_10():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, color=WHITE)

    # Centered layout with generous whitespace
    center_x = Inches(3.0)
    center_w = Inches(7.333)

    # "Bewo" logo text
    add_text(slide, center_x, Inches(1.8), center_w, Inches(1.0),
             "Bewo", font_size=64, color=INK, bold=True, font_name=SERIF, alignment=PP_ALIGN.CENTER)

    # Tagline with mixed styling
    txbox = slide.shapes.add_textbox(center_x, Inches(2.9), center_w, Inches(0.8))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "Predicting health crises "
    r1.font.size = Pt(20)
    r1.font.color.rgb = SECONDARY
    r1.font.name = SERIF
    r2 = p.add_run()
    r2.text = "before"
    r2.font.size = Pt(20)
    r2.font.color.rgb = INK
    r2.font.bold = True
    r2.font.name = SERIF
    r3 = p.add_run()
    r3.text = " they happen."
    r3.font.size = Pt(20)
    r3.font.color.rgb = SECONDARY
    r3.font.name = SERIF

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    r4 = p2.add_run()
    r4.text = "For Singapore's 440,000 diabetics."
    r4.font.size = Pt(16)
    r4.font.color.rgb = MUTED
    r4.font.name = SERIF

    # Divider
    add_divider(slide, Inches(5.0), Inches(4.0), Inches(3.333), color=HAIRLINE)

    # Stats row
    stat_items = [
        ("18", "AI TOOLS"),
        ("54", "API ROUTES"),
        ("20,000", "LINES OF CODE"),
        ("3", "LIVE VIEWS"),
    ]

    for i, (num, label) in enumerate(stat_items):
        sx = Inches(2.8) + Inches(i * 2.1)
        add_text(slide, sx, Inches(4.4), Inches(1.8), Inches(0.4),
                 num, font_size=26, color=NAVY, bold=True, font_name=MONO, alignment=PP_ALIGN.CENTER)
        add_text(slide, sx, Inches(4.85), Inches(1.8), Inches(0.25),
                 label, font_size=8, color=MUTED, font_name=MONO, alignment=PP_ALIGN.CENTER)

    # Live prototype badge
    badge_w = Inches(3.0)
    badge_x = Inches(5.167)
    badge_y = Inches(5.5)
    add_shape(slide, badge_x, badge_y, badge_w, Inches(0.35),
              fill_color=RGBColor(0xF0, 0xFD, 0xF4), corner_radius=True)
    add_text(slide, badge_x, badge_y + Inches(0.03), badge_w, Inches(0.3),
             "\u25cf  Live prototype — not a mockup", font_size=11, color=EMERALD, bold=True,
             font_name=MONO, alignment=PP_ALIGN.CENTER)

    # "Before crisis. Not after."
    add_text(slide, center_x, Inches(6.1), center_w, Inches(0.4),
             "Before crisis. Not after.", font_size=18, color=INK, bold=True,
             font_name=SERIF, alignment=PP_ALIGN.CENTER)

    # Footer
    add_text(slide, MARGIN_L, Inches(6.8), Inches(11.333), Inches(0.3),
             "NUS x Synapxe x IMDA  |  AI Innovation Challenge 2026  |  Bewo Health",
             font_size=9, color=MUTED, font_name=MONO, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ═══════════════════════════════════════════════════════════

build_slide_1()
build_slide_2()
build_slide_3()
build_slide_4()
build_slide_5()
build_slide_6()
build_slide_7()
build_slide_8()
build_slide_9()
build_slide_10()

output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "Bewo_Pitch_Deck.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
